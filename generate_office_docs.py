from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PPTXInches
import os

def create_docx(md_path, out_path):
    doc = Document()
    with open(md_path, 'r', encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if not line: continue
            if line.startswith('# '):
                doc.add_heading(line[2:], level=1)
            elif line.startswith('## '):
                doc.add_heading(line[3:], level=2)
            elif line.startswith('### '):
                doc.add_heading(line[4:], level=3)
            elif line.startswith('- '):
                doc.add_paragraph(line[2:], style='List Bullet')
            elif line[0].isdigit() and line[1:3] == '. ':
                doc.add_paragraph(line[3:], style='List Number')
            else:
                doc.add_paragraph(line.replace('**', '').replace('*', ''))
                
    # Add generated diagrams at the end of the proposal document for reference
    doc.add_page_break()
    doc.add_heading('Appendix: System Architecture Diagrams', level=1)
    
    doc.add_heading('4-Agent SME Toolkit Architecture', level=2)
    if os.path.exists('arch_flow.png'):
        doc.add_picture('arch_flow.png', width=Inches(6.0))
        
    doc.add_heading('Key Feature Venn Diagram Method', level=2)
    if os.path.exists('venn_flow.png'):
        doc.add_picture('venn_flow.png', width=Inches(6.0))
        
    doc.add_heading('Complete Technical Foundation', level=2)
    if os.path.exists('db_flow.png'):
        doc.add_picture('db_flow.png', width=Inches(6.0))
        
    doc.save(out_path)

def create_pptx(md_path, out_path):
    prs = Presentation()
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    slides = content.split('<!-- slide -->')
    for slide_text in slides:
        lines = [L for L in slide_text.splitlines() if L.strip()]
        if not lines: continue
        
        title = ""
        bullets = []
        is_mermaid = False
        
        for line in lines:
            stripped = line.strip()
            if stripped.startswith('````') or stripped.startswith('```'):
                is_mermaid = not is_mermaid
                continue
            if is_mermaid: 
                continue # Skip mermaid code block
                
            if stripped.startswith('# '): 
                title = stripped[2:].replace('*', '').replace('_', '')
            elif stripped.startswith('## '): 
                if not title: 
                    title = stripped[3:].replace('*', '').replace('_', '')
                else: 
                    bullets.append(stripped)
            else:
                if stripped != '---' and stripped != '----':
                    bullets.append(line)
                    
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # title and content layout
        if slide.shapes.title and title:
            slide.shapes.title.text = title
            
        if len(slide.shapes.placeholders) > 1 and bullets:
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            for b in bullets:
                stripped_b = b.strip()
                p = tf.add_paragraph()
                clean_text = stripped_b.replace('**', '').replace('*', '').replace('## ', '').replace('### ', '').replace('<br/>', ' ')
                
                if stripped_b.startswith('- '):
                    p.level = 0
                    p.text = clean_text[2:]
                elif b.startswith('  - '):
                    p.level = 1
                    p.text = clean_text[2:]
                elif stripped_b[0].isdigit() and stripped_b[1:3] == '. ':
                    p.level = 0
                    p.text = clean_text
                else:
                    p.level = 0
                    p.text = clean_text
        
        # Insert image based on slide title matching
        if "4-Agent SME Toolkit Architecture" in title and os.path.exists('arch_flow.png'):
            # Clear text frame 1 and add picture
            if slide.shapes.placeholders:
                sp = slide.shapes.placeholders[1]
                sp.element.getparent().remove(sp.element)
            prs.slides[-1].shapes.add_picture('arch_flow.png', PPTXInches(1.0), PPTXInches(2.0), width=PPTXInches(8.0))
            
        elif "Diffusion Agent: Distractor Modeling" in title and os.path.exists('venn_flow.png'):
            # This slide contains some bullets so we'll offset the picture to the right or bottom.
            prs.slides[-1].shapes.add_picture('venn_flow.png', PPTXInches(1.5), PPTXInches(4.5), width=PPTXInches(7.0))
            
        elif "Complete Technical Foundation" in title and os.path.exists('db_flow.png'):
             if slide.shapes.placeholders:
                try:
                    sp = slide.shapes.placeholders[1]
                    sp.element.getparent().remove(sp.element)
                except KeyError:
                    pass
             prs.slides[-1].shapes.add_picture('db_flow.png', PPTXInches(1.0), PPTXInches(2.0), width=PPTXInches(8.0))
            
    prs.save(out_path)

prop_md = "/Users/jakecho/.gemini/antigravity/brain/18d55a34-aeb6-45f9-a8d6-20b7e61615b5/ice_conference_proposal.md"
pres_md = "/Users/jakecho/.gemini/antigravity/brain/18d55a34-aeb6-45f9-a8d6-20b7e61615b5/ice_conference_presentation.md"

create_docx(prop_md, "ICE_Conference_Proposal_Updated.docx")
create_pptx(pres_md, "ICE_Conference_Presentation_Updated.pptx")
print("Successfully generated updated documents!")
